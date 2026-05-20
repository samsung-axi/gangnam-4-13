import os
from datetime import datetime
from typing import List, Optional
from uuid import UUID
import aiofiles
import aioboto3
import PyPDF2
import docx
import openpyxl
from botocore.exceptions import ClientError
from dotenv import load_dotenv
from fastapi import UploadFile, HTTPException
from sentence_transformers import SentenceTransformer
from openai import AsyncOpenAI
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession
from sqlalchemy.orm import sessionmaker
from sqlalchemy.future import select
import base64
import fitz
import pptx

from app.models.interdoc import Interdoc

load_dotenv()

DATABASE_URL = os.getenv('CONNECTION_STRING').replace('postgresql://', 'postgresql+asyncpg://')
engine = create_async_engine(DATABASE_URL, echo=True)
AsyncSessionLocal = sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

async def get_db():
    async with AsyncSessionLocal() as session:
        try:
            yield session
        finally:
            await session.close()

# S3 클라이언트 설정
session = aioboto3.Session(
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY'),
    aws_secret_access_key=os.getenv('AWS_SECRET_KEY'),
    region_name=os.getenv('AWS_REGION')
)

# OpenAI 클라이언트
openai_client = AsyncOpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# 문서 임베딩 모델 초기화
model = SentenceTransformer('sentence-transformers/distiluse-base-multilingual-cased-v2')

async def read_file_content(file: UploadFile) -> str:
    """파일 형식에 따라 내용을 읽는 함수"""
    content = ""
    file_ext = file.filename.lower().split('.')[-1]
    
    await file.seek(0) # 중요: 파일 포인터 초기화
    
    try:
        if file_ext == 'txt':
            content = (await file.read()).decode('utf-8')
        
        elif file_ext == 'pdf':
            # 임시 파일로 저장 후 처리
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False) as temp_file:
                await temp_file.write(await file.read())
                temp_path = temp_file.name

            try:
                pdf_reader = PyPDF2.PdfReader(temp_path)
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
        
        elif file_ext in ['doc', 'docx']:
            # 임시 파일로 저장 후 처리
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}') as temp_file:
                await temp_file.write(await file.read())
                temp_path = temp_file.name

            try:
                doc = docx.Document(temp_path)
                # 본문 텍스트
                texts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
                # 표 텍스트
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            texts.append(' | '.join(row_text))
                content = "\n".join(texts)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            
        elif file_ext in ['xlsx', 'xls']:
            # 임시 파일로 저장 후 처리
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}') as temp_file:
                await temp_file.write(await file.read())
                temp_path = temp_file.name
            
            try:
                workbook = openpyxl.load_workbook(temp_path, data_only=True)
                sheets_content = []
                
                for sheet in workbook.sheetnames:
                    worksheet = workbook[sheet]
                    sheet_content = []
                    sheet_content.append(f"[시트: {sheet}]")
                    for row in worksheet.iter_rows():
                        row_values = [str(cell.value) if cell.value is not None else '' for cell in row]
                        if any(row_values):
                            sheet_content.append(' | '.join(row_values))
                    
                    sheets_content.append('\n'.join(sheet_content))
                
                content = '\n\n'.join(sheets_content)
                
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
        
        elif file_ext in ['ppt', 'pptx']:
            # 임시 파일로 저장 후 처리
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}') as temp_file:
                await temp_file.write(await file.read())
                temp_path = temp_file.name
            try:
                prs = pptx.Presentation(temp_path)
                slides_text = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        slides_text.append("\n".join(slide_text))
                content = "\n\n".join(slides_text)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
        
        else:
            raise HTTPException(
                status_code=400,
                detail="지원하지 않는 파일 형식입니다. (지원 형식: txt, pdf, doc, docx, xlsx, xls, ppt, pptx)"
            )
            
        if not content.strip():
            raise HTTPException(
                status_code=400,
                detail="파일에서 텍스트를 추출할 수 없습니다."
            )
            
        return content
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"파일 읽기 실패: {str(e)}"
        )

async def pdf_to_images_base64(pdf_path):
    doc = fitz.open(pdf_path)
    images_base64 = []
    for page in doc:
        pix = page.get_pixmap()
        img_bytes = pix.tobytes("png")
        img_b64 = base64.b64encode(img_bytes).decode("utf-8")
        images_base64.append(img_b64)
    return images_base64

async def extract_text_from_file(file: UploadFile) -> str:
    """파일명과 내용을 포함한 요약 생성 함수"""
    try:
        await file.seek(0)
        content = ""
        filename = file.filename
        file_ext = filename.lower().split('.')[-1]
        is_pdf = file_ext == 'pdf'
        try:
            content = await read_file_content(file)
        except HTTPException as e:
            if is_pdf and (e.status_code == 400 or e.status_code == 500):
                # PDF 텍스트 추출 실패 시 이미지 변환
                async with aiofiles.tempfile.NamedTemporaryFile(delete=False) as temp_file:
                    await file.seek(0)
                    await temp_file.write(await file.read())
                    temp_path = temp_file.name
                try:
                    content = await pdf_to_images_base64(temp_path)
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
            else:
                raise e
        # content가 비어있으면(pdf에서 텍스트 추출 실패 등) 이미지 변환
        if is_pdf and (not content or (isinstance(content, str) and not content.strip())):
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False) as temp_file:
                await file.seek(0)
                await temp_file.write(await file.read())
                temp_path = temp_file.name
            try:
                content = await pdf_to_images_base64(temp_path)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
        # content가 이미지 리스트라면 vision 프롬프트로 전달
        if isinstance(content, list):
            prompt = f"""
파일명: {filename}

아래 이미지는 PDF 각 페이지를 이미지로 변환한 것입니다. 이미지를 보고 문서의 용도와 종류를 한 문장으로 요약해줘.

예시: '프로젝트 기획안 작성을 위한 문서'
조건:
- 문서의 구성요소(예: 목적, 일정, 예산 등)를 기준으로 문서의 종류를 추론할 것
- 문서 형식이나 스타일 설명은 제외
- 구체적인 내용이 아니라, '무슨 목적으로 어떤 형식으로 작성된 문서'인지 설명할 것
- 반드시 단문형 서술체로 요약할 것 (예: '프로젝트 기획안 작성을 위한 문서')
"""
            messages = [
                {"role": "system", "content": "당신은 문서 분석 전문가입니다. 이미지로 된 문서의 용도와 종류를 한 문장으로 요약하세요."},
                {"role": "user", "content": prompt.strip()}
            ]
            for img_b64 in content[:2]:
                messages.append({
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                    ]
                })
            response = await openai_client.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                max_tokens=300,
                temperature=0.3
            )
            summary = response.choices[0].message.content.strip()
            return summary
        # 텍스트라면 기존 프롬프트로 처리
        prompt = f"""
파일명: {filename}

문서 내용:
{content}

위 문서는 어떤 목적이나 업무에 사용되는 문서인지 판단해서 아래 형식처럼 한 문장으로 요약해줘.

예시: "프로젝트 기획안 작성을 위한 문서"
조건:
- 문서의 구성요소(예: 목적, 일정, 예산 등)를 기준으로 문서의 종류를 추론할 것
- 문서 형식이나 스타일 설명은 제외
- 구체적인 내용이 아니라, '무슨 목적으로 어떤 형식으로 작성된 문서'인지 설명할 것
- 반드시 단문형 서술체로 요약할 것 (예: "프로젝트 기획안 작성을 위한 문서")
"""
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "당신은 문서 분석 전문가입니다. 문서의 제목과 내용을 바탕으로 문서의 용도와 종류를 한 문장으로 요약하세요."},
                {"role": "user", "content": prompt.strip()}
            ],
            max_tokens=300,
            temperature=0.3
        )
        summary = response.choices[0].message.content.strip()
        return summary
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"문서 요약 실패: {str(e)}"
        )

async def create_document(
    db: AsyncSession,
    file: UploadFile,
    doc_type: str,
    update_user_id: UUID
) -> Interdoc:
    """문서 생성 함수"""
    try:
        # 파일명에서 '예스폼_' 접두사 제거
        filename = file.filename
        if filename.startswith('예스폼_'):
            filename = filename[len('예스폼_'):]
        file.filename = filename
        content = await extract_text_from_file(file)
        embedding = model.encode(content)
        
        s3_path = f"documents/{filename}"
        
        await file.seek(0)
        async with session.client('s3') as s3:
            await s3.upload_fileobj(
                file.file,
                os.getenv('AWS_BUCKET_NAME'),
                s3_path
            )
        
        doc = Interdoc(
            interdocs_type_name=doc_type,
            interdocs_filename=filename,
            interdocs_contents=content[:255],
            interdocs_vector=embedding,
            interdocs_path=s3_path,
            interdocs_uploaded_date=datetime.now(),
            interdocs_update_user_id=update_user_id
        )
        
        db.add(doc)
        await db.commit()
        await db.refresh(doc)
        
        return doc
        
    except Exception as e:
        if 's3_path' in locals():
            try:
                async with session.client('s3') as s3:
                    await s3.delete_object(
                        Bucket=os.getenv('AWS_BUCKET_NAME'),
                        Key=s3_path
                    )
            except:
                pass
        await db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

async def update_document(
    db: AsyncSession,
    doc_id: UUID,
    file: UploadFile,
    update_user_id: UUID
) -> Interdoc:
    """문서 수정 함수"""
    try:
        query = select(Interdoc).filter(Interdoc.interdocs_id == doc_id)
        result = await db.execute(query)
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
            
        content = await extract_text_from_file(file)
        embedding = model.encode(content)
        
        old_s3_path = doc.interdocs_path
        new_s3_path = f"documents/{file.filename}"
        
        await file.seek(0)
        async with session.client('s3') as s3:
            await s3.upload_fileobj(
                file.file,
                os.getenv('AWS_BUCKET_NAME'),
                new_s3_path
            )
            
            await s3.delete_object(
                Bucket=os.getenv('AWS_BUCKET_NAME'),
                Key=old_s3_path
            )
        
        doc.interdocs_filename = file.filename
        doc.interdocs_contents = content[:255]
        doc.interdocs_vector = embedding
        doc.interdocs_path = new_s3_path
        doc.interdocs_updated_date = datetime.now()
        doc.interdocs_update_user_id = update_user_id
        
        await db.commit()
        await db.refresh(doc)
        
        return doc
        
    except Exception as e:
        if 'new_s3_path' in locals():
            try:
                async with session.client('s3') as s3:
                    await s3.delete_object(
                        Bucket=os.getenv('AWS_BUCKET_NAME'),
                        Key=new_s3_path
                    )
            except:
                pass
        await db.rollback()
        raise HTTPException(status_code=500, detail=str(e))

async def get_documents(
    db: AsyncSession,
    skip: int = 0,
    limit: int = 200
) -> List[Interdoc]:
    """문서 목록 조회 함수"""
    print("get_documents 함수 시작")
    try:
        print(f"DB 세션 객체: {db}")
        print(f"Interdoc 모델: {Interdoc}")
        
        query = select(Interdoc).offset(skip).limit(limit)
        result = await db.execute(query)
        documents = result.scalars().all()
        
        print(f"조회된 문서 수: {len(documents)}")
        return documents
    except Exception as e:
        print(f"!!! get_documents 에러 발생: {e}")
        raise HTTPException(status_code=500, detail=f"문서 목록 조회 실패: {str(e)}")

async def get_document(
    db: AsyncSession,
    doc_id: UUID
) -> Optional[Interdoc]:
    """단일 문서 조회 함수"""
    try:
        query = select(Interdoc).filter(Interdoc.interdocs_id == doc_id)
        result = await db.execute(query)
        return result.scalar_one_or_none()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"단일 문서 조회 실패: {str(e)}")

async def delete_document(
    db: AsyncSession,
    doc_id: UUID
) -> bool:
    """문서 삭제 함수"""
    try:
        query = select(Interdoc).filter(Interdoc.interdocs_id == doc_id)
        result = await db.execute(query)
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
            
        async with session.client('s3') as s3:
            await s3.delete_object(
                Bucket=os.getenv('AWS_BUCKET_NAME'),
                Key=doc.interdocs_path
            )
        
        await db.delete(doc)
        await db.commit()
        
        return True
        
    except ClientError as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"S3 삭제 실패: {str(e)}")
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=str(e))