"""SPOTTER ABM × DB 브리핑 PPT — 실측 지표 기반.

기준: origin/dev (alembic a8f3d2e7c1b9), 2026-05-06
산출: docs/presentation/spotter-abm-db-briefing.pptx

실 코드 + DB 쿼리 검증:
- DB 총 5.3 GB, 87 public 테이블, 213 인덱스, 43 FK, 1,153 컬럼
- ABM /simulate-abm: n_agents=5000, default policy mode (LLM 0회)
- RAG 임베딩 10,255건 (법률 57건 + 판례 222건)

실행: python -m scripts.build_abm_db_briefing_pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# 색
# ---------------------------------------------------------------------------
BG = RGBColor(0x10, 0x14, 0x20)
PANEL = RGBColor(0x1A, 0x22, 0x36)
PANEL_ALT = RGBColor(0x12, 0x18, 0x28)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_2 = RGBColor(0xFF, 0xB7, 0x4D)
TEXT = RGBColor(0xE3, 0xE9, 0xF4)
MUTED = RGBColor(0x9C, 0xA9, 0xBE)
GREEN = RGBColor(0x6E, 0xE7, 0xB7)
RED = RGBColor(0xFF, 0x9A, 0x9A)
PURPLE = RGBColor(0xC4, 0x9B, 0xFF)


# ---------------------------------------------------------------------------
# 헬퍼
# ---------------------------------------------------------------------------
def add_dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"


def add_bullets(slide, x, y, w, h, lines, *, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, c = item
        else:
            txt, c = item, TEXT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"·  {txt}"
        run.font.size = Pt(size)
        run.font.color.rgb = c
        run.font.name = "Pretendard"


def add_table(slide, x, y, w, headers, rows, *, header_size=11, body_size=10, row_h=0.38):
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(row_h * (len(rows) + 1) + 0.1),
    ).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = ACCENT
                r.font.name = "Pretendard"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            if isinstance(val, tuple):
                txt, color = val
            else:
                txt, color = val, TEXT
            cell.text = str(txt)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = color
                    run.font.name = "Pretendard"
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_ALT


def add_band(slide, y, label, *, color=ACCENT):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.04))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    add_text(slide, 0.5, y + 0.05, 12, 0.4, label, size=14, bold=True, color=color)


def add_metric_card(slide, x, y, w, h, label, value, *, value_color=ACCENT):
    """KPI 카드 — 라벨 + 큰 숫자."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = PANEL
    card.shadow.inherit = False
    add_text(slide, x + 0.15, y + 0.08, w - 0.3, 0.35, label, size=11, color=MUTED)
    add_text(slide, x + 0.15, y + 0.42, w - 0.3, 0.7, value, size=22, bold=True, color=value_color)


# ---------------------------------------------------------------------------
# Slide 1 — 표지
# ---------------------------------------------------------------------------
def build_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.7, 1.5, 12, 1.0, "SPOTTER", size=56, bold=True, color=ACCENT)
    add_text(s, 0.7, 2.6, 12, 0.8, "ABM × DB 인프라 브리핑", size=32, bold=True, color=TEXT)
    add_text(
        s,
        0.7,
        3.5,
        12,
        0.5,
        "기준 시점: 2026-05-06 · origin/dev · alembic a8f3d2e7c1b9",
        size=15,
        color=MUTED,
    )
    add_band(s, 4.5, "구성")
    add_bullets(
        s,
        0.7,
        4.95,
        12,
        2.5,
        [
            ("Slide 2 — ABM 시스템 개요 (5,000 에이전트 × 3-Tier × 4 mode)", TEXT),
            ("Slide 3 — ABM 의사결정 흐름 + 메모리 시스템 (활성 / 비활성)", TEXT),
            ("Slide 4 — DB 인프라 (87 테이블 / 5.3 GB / 213 인덱스)", TEXT),
            ("Slide 5 — 데이터 자산 (마포 1,070만 / FTC 35K brand / RAG 10K 임베딩)", TEXT),
            ("Slide 6 — 통합 흐름 + 정량 요약", TEXT),
        ],
        size=14,
    )


# ---------------------------------------------------------------------------
# Slide 2 — ABM 개요
# ---------------------------------------------------------------------------
def build_abm_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "ABM 시스템 개요", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "마포구 5,000 가상 에이전트가 매 시간 행동 결정 → 매장별 매출 추정",
        size=13,
        color=MUTED,
    )

    # KPI 카드 4 개
    add_metric_card(s, 0.5, 1.5, 3.0, 1.2, "에이전트 수", "5,000")
    add_metric_card(s, 3.7, 1.5, 3.0, 1.2, "Tier", "S 50 / A 200 / B 4,750")
    add_metric_card(s, 6.9, 1.5, 3.0, 1.2, "기본 LLM 호출", "0", value_color=GREEN)
    add_metric_card(s, 10.1, 1.5, 2.7, 1.2, "Mode", "4")

    # 인구 구성
    add_band(s, 3.0, "Population Mix (마포 실측 비례, n_personas 비례 scale)")
    add_table(
        s,
        0.5,
        3.45,
        12.3,
        ["역할", "비율", "근거", "행동 패턴"],
        [
            ["residents (거주민)", "60%", "SGIS 361,380", "일상 생활 — 식사·카페·소비"],
            ["ext_commuters (외부 통근)", "20%", "마포 사업체 종사 281,385 일부", "출근 시간 진입 / 퇴근 이탈"],
            ["commuters (마포 내 통근)", "10%", "거주+근무 일치", "거주 동 ↔ 근무 동"],
            ["ext_visitors (외부 방문)", "5%", "홍대·연남 야간 방문", "저녁 시간 진입"],
            ["visitors (단기 방문)", "4%", "마포 내 단기 이동", "비정기 방문"],
            ["owners (점주)", "1%", "운영 매장 보유", "9~22시 매장 상주"],
        ],
        body_size=10,
    )

    # 한 줄 요약 박스
    add_band(s, 6.4, "핵심 한 줄", color=ACCENT_2)
    add_text(
        s,
        0.5,
        6.85,
        12.3,
        0.5,
        "에이전트 = 통계 분포 + 정책 기반 행동 + (옵션) Tier S 50명 LLM 추론",
        size=14,
        color=GREEN,
    )


# ---------------------------------------------------------------------------
# Slide 3 — ABM 의사결정 + 메모리
# ---------------------------------------------------------------------------
def build_abm_decision(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "ABM 의사결정 흐름", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "agents.py:decide() 라우터 — 4 mode 분기. /simulate-abm 기본 = use_policy=True",
        size=13,
        color=MUTED,
    )

    # 4 mode 표
    add_band(s, 1.4, "4가지 Mode")
    add_table(
        s,
        0.5,
        1.85,
        12.3,
        ["Mode", "Tier S (50)", "Tier A (200)", "Tier B (≈4,750)", "LLM 호출"],
        [
            [
                ("default (use_policy=True)", GREEN),
                "policy",
                "policy",
                "policy",
                ("0", GREEN),
            ],
            [
                "enable_llm_decisions",
                ("smart_decide (LLM)", ACCENT_2),
                "policy",
                "policy",
                "Tier S 만 (≈1,200/sim)",
            ],
            [
                "풀 JSON 모드",
                "smart_decide",
                ("fast_decide (LLM)", ACCENT_2),
                "rule",
                "S+A",
            ],
            ["DSL 모드", "dsl_decide", "dsl_decide", "dsl_decide", ("전 Tier", RED)],
        ],
        row_h=0.42,
    )

    # 핵심 트릭
    add_band(s, 4.0, "비용 최소화 트릭")
    add_bullets(
        s,
        0.5,
        4.45,
        6.4,
        2.5,
        [
            ("policy_generator — Ollama Qwen2.5:3b 로컬", TEXT),
            ("  → role × 날씨 × 시간대 = 11개 정책 cold gen", MUTED),
            ("  → policy_cache.json 저장 → 매 시뮬 LLM 0회", GREEN),
            ("Archetype 30+ multiplier", TEXT),
            ("  → resident 7 / commuter 5 / visitor 4 ...", MUTED),
            ("  → 같은 role 안 행동 다양성 확보", MUTED),
        ],
        size=12,
    )

    # 메모리 시스템 표
    add_text(s, 7.0, 4.4, 6, 0.4, "메모리 / 보조 시스템 (현재 활성)", size=13, bold=True, color=ACCENT_2)
    add_table(
        s,
        7.0,
        4.85,
        6.0,
        ["시스템", "활성"],
        [
            ["MemoryStore (raw + 일일 요약)", ("✅ 사용", GREEN)],
            ["policy_cache.json (LLM 0회 정책)", ("✅ 사용", GREEN)],
            ["Archetype 30+ multiplier", ("✅ 사용", GREEN)],
            ["Memory Seeder (14일 prefill)", ("❌ off", RED)],
            ["PgVectorMemory (semantic search)", ("❌ off", RED)],
        ],
        body_size=11,
        row_h=0.36,
    )


# ---------------------------------------------------------------------------
# Slide 4 — DB 인프라
# ---------------------------------------------------------------------------
def build_db_infra(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "DB 인프라", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "PostgreSQL (RDS) · SQLAlchemy ORM · Alembic 마이그레이션",
        size=13,
        color=MUTED,
    )

    # KPI 5개
    add_metric_card(s, 0.5, 1.5, 2.4, 1.2, "DB 총 크기", "5.3 GB")
    add_metric_card(s, 3.0, 1.5, 2.4, 1.2, "Public 테이블", "87")
    add_metric_card(s, 5.5, 1.5, 2.4, 1.2, "ORM 클래스", "77")
    add_metric_card(s, 8.0, 1.5, 2.4, 1.2, "Alembic 마이그", "37")
    add_metric_card(s, 10.5, 1.5, 2.3, 1.2, "FK / 인덱스", "43 / 213")

    # Top 10 size 표
    add_band(s, 3.0, "TOP 10 테이블 (디스크 사용량)")
    add_table(
        s,
        0.5,
        3.45,
        6.4,
        ["테이블", "크기"],
        [
            ["living_population_grid", "2.96 GB"],
            ["bus_boarding_daily", "640 MB"],
            ["living_population", "623 MB"],
            ["langchain_pg_embedding (RAG)", "276 MB"],
            ["seoul_adstrd_stor", "170 MB"],
            ["golmok_commercial", "128 MB"],
            ["seoul_ttareungi_usage_daily", "89 MB"],
            ["jeonse_monthly_rent", "76 MB"],
            ["district_sales_seoul", "71 MB"],
            ["seoul_district_sales", "48 MB"],
        ],
        row_h=0.32,
        body_size=10,
    )

    # 마이그레이션 / 운영 안정성
    add_text(s, 7.2, 3.05, 6, 0.4, "운영 안정성", size=13, bold=True, color=ACCENT_2)
    add_bullets(
        s,
        7.2,
        3.5,
        6.0,
        3.5,
        [
            ("alembic head 정합 (phantom revision 복구 후)", GREEN),
            ("SQLAlchemy Engine 싱글톤화 → RDS 포화 해소", GREEN),
            ("vector_db pool_recycle=1800 (idle 회수)", GREEN),
            ("dong_code FK 4 그룹 통합 (425 + 399 master)", GREEN),
            ("87 테이블 ORM ↔ DB drift 정합 (zombie 정리)", GREEN),
            ("외부 API NULL fill (subway / ttareungi / hotspots)", GREEN),
            ("ETL 재적재: ttareungi.dong_code 마포 100%", GREEN),
            ("권한: users.is_superadmin BOOLEAN 컬럼 신설", GREEN),
        ],
        size=11,
    )


# ---------------------------------------------------------------------------
# Slide 5 — 데이터 자산
# ---------------------------------------------------------------------------
def build_data_assets(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "데이터 자산 (현재 적재)", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "외부 API 8종 적재 + 마포 한정 데이터셋 + 법률 RAG",
        size=13,
        color=MUTED,
    )

    # 핵심 row count
    add_band(s, 1.4, "주요 테이블 row count")
    add_table(
        s,
        0.5,
        1.85,
        6.4,
        ["테이블", "row 수"],
        [
            ["living_population_grid (마포)", "10,543,528"],
            ["living_population (마포)", "968,064"],
            ["seoul_adstrd_stor (마포)", "101,871"],
            ["kakao_store_menu", "81,037"],
            ["ftc_brand_franchise", "34,708"],
            ["seoul_district_sales", "87,938"],
            ["langchain_pg_embedding (RAG)", "10,255"],
            ["biz_brand_mapping", "5,867"],
            ["kakao_store", "4,418"],
            ["weather_daily", "2,665"],
        ],
        row_h=0.32,
        body_size=10,
    )

    # 외부 API 출처
    add_text(s, 7.2, 1.4, 6, 0.4, "외부 API 출처 (8종)", size=13, bold=True, color=ACCENT_2)
    add_bullets(
        s,
        7.2,
        1.85,
        6.0,
        2.7,
        [
            ("서울열린데이터광장 (flpop / golmok / 지하철)", TEXT),
            ("SGIS — 인구 / 가구 / 사업체 (2026 KOSIS)", TEXT),
            ("MOLIT — 부동산 실거래", TEXT),
            ("공정거래위원회 (FTC) — 프랜차이즈 정보공개서", TEXT),
            ("ECOS — 한국은행 경기 지표 (cycle 100%)", TEXT),
            ("기상청 — weather_daily", TEXT),
            ("Kakao Local Search — kakao_store / menu", TEXT),
            ("공공자전거 — ttareungi 마포 dong_code 100%", TEXT),
        ],
        size=11,
    )

    # 법률 RAG
    add_band(s, 4.7, "법률 RAG (specialist 4 + 판례 인용)")
    add_table(
        s,
        0.5,
        5.15,
        12.3,
        ["항목", "값"],
        [
            ["임베딩 (langchain_pg_embedding)", "10,255 vectors (BGE-M3)"],
            ["법률 조항 (law_legislations)", "57"],
            ["판례 (law_precedents)", "222"],
            ["specialist 4종", "가맹사업법 · 공정거래법 · 식품위생법 · 건축법"],
        ],
        row_h=0.36,
        body_size=11,
    )


# ---------------------------------------------------------------------------
# Slide 6 — 통합 흐름 + 정량 요약
# ---------------------------------------------------------------------------
def build_integration(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "ABM ↔ DB 통합 흐름", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "DB → world_loader → 5,000 agents → 매장 선택 → 매출 집계 → 응답",
        size=13,
        color=MUTED,
    )

    # 흐름 다이어그램 (텍스트 박스 5개 가로 정렬)
    flow = [
        ("DB", "87 테이블\n5.3 GB", PURPLE),
        ("world_loader", "stores / dongs\n로드", ACCENT),
        ("agents 5,000", "Tier S/A/B\n분배", ACCENT),
        ("decide()", "policy\n(LLM 0회)", GREEN),
        ("응답", "trajectory +\n매출 집계", ACCENT_2),
    ]
    box_w = 2.3
    gap = 0.2
    start_x = 0.5
    for i, (title, sub, color) in enumerate(flow):
        x = start_x + i * (box_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(box_w), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.shadow.inherit = False
        add_text(s, x + 0.1, 1.7, box_w - 0.2, 0.4, title, size=14, bold=True, color=color)
        add_text(s, x + 0.1, 2.15, box_w - 0.2, 0.7, sub, size=11, color=TEXT)
        # arrow
        if i < len(flow) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_w + 0.01),
                Inches(2.15),
                Inches(gap - 0.02),
                Inches(0.2),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    # 정량 요약
    add_band(s, 3.2, "정량 요약 (production 기준)")
    add_table(
        s,
        0.5,
        3.65,
        6.2,
        ["ABM 지표", "값"],
        [
            ["에이전트 수", "5,000"],
            ["LLM 호출 (default mode)", ("0", GREEN)],
            ["LLM 호출 (Tier S 모드, 24h)", "≈1,200"],
            ["Population mix base 합", "500 (n_agents 비례 scale)"],
            ["Tier S 50 + A 200 + B ≈4,750", "(고정/고정/잔여)"],
        ],
        body_size=11,
    )
    add_table(
        s,
        7.0,
        3.65,
        6.0,
        ["DB 지표", "값"],
        [
            ["DB 총 크기", "5.3 GB"],
            ["테이블 / ORM / 마이그", "87 / 77 / 37"],
            ["FK / 인덱스 / 컬럼", "43 / 213 / 1,153"],
            ["RAG 임베딩", "10,255 (BGE-M3)"],
            ["users / superadmin", "23 / 1"],
        ],
        body_size=11,
    )

    # 마무리 메시지
    add_band(s, 6.3, "한 줄 결론", color=ACCENT_2)
    add_text(
        s,
        0.5,
        6.75,
        12.3,
        0.5,
        "ABM 5,000 에이전트 시뮬을 매 호출 LLM 0회로 돌릴 수 있도록 정책 캐시 + 메모리 + DB 인프라가 정합 상태",
        size=13,
        color=GREEN,
    )


# ---------------------------------------------------------------------------
def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title(prs)
    build_abm_overview(prs)
    build_abm_decision(prs)
    build_db_infra(prs)
    build_data_assets(prs)
    build_integration(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "spotter-abm-db-briefing.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"saved: {out}")
    return out


if __name__ == "__main__":
    main()
