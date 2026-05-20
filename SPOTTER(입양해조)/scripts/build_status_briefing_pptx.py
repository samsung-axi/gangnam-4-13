"""SPOTTER 현재 상태 브리핑 PPT — DB / 코드 실측 기반.

산출: docs/presentation/spotter-status-briefing.pptx

기준 시점: 2026-05-06 (origin/dev HEAD f5ee308c, alembic a8f3d2e7c1b9)
실행: python -m scripts.build_status_briefing_pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

BG = RGBColor(0x10, 0x14, 0x20)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_2 = RGBColor(0xFF, 0xB7, 0x4D)
TEXT = RGBColor(0xE3, 0xE9, 0xF4)
MUTED = RGBColor(0x9C, 0xA9, 0xBE)
GREEN = RGBColor(0x6E, 0xE7, 0xB7)
RED = RGBColor(0xFF, 0x9A, 0x9A)


def add_dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"


def add_bullets(slide, x, y, w, h, lines, *, size=13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, c = item
        else:
            txt, c = item, TEXT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"·  {txt}"
        run.font.size = Pt(size)
        run.font.color.rgb = c
        run.font.name = "Pretendard"


def add_kv_table(slide, x, y, w, headers, rows, *, header_size=12, body_size=11):
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.38 * (len(rows) + 1) + 0.1),
    ).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = ACCENT
                r.font.name = "Pretendard"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x22, 0x36)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = TEXT
                    run.font.name = "Pretendard"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x28)


def add_band(slide, y, label):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.04))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    add_text(slide, 0.5, y + 0.05, 12, 0.4, label, size=14, bold=True, color=ACCENT)


def build_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.7, 1.6, 12, 1.0, "SPOTTER — 현재 상태 브리핑", size=42, bold=True)
    add_text(
        s,
        0.7,
        2.7,
        12,
        0.5,
        "기준: origin/dev (alembic a8f3d2e7c1b9) · 2026-05-06",
        size=18,
        color=MUTED,
    )
    add_band(s, 4.0, "오늘 다룰 영역")
    add_bullets(
        s,
        0.7,
        4.5,
        12,
        2.5,
        [
            ("ABM 시뮬레이션 — 5,000 agents 매장별 매출 추정", TEXT),
            ("DB / 데이터 — 87 테이블 / 외부 API 통합", TEXT),
            ("권한 시스템 — master / manager / superadmin 3계층", TEXT),
        ],
        size=16,
    )


def build_abm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "ABM 시뮬레이션 (현재)", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "마포구 5,000 가상 에이전트가 매 시간 행동 결정 → 매장별 매출 추정",
        size=13,
        color=MUTED,
    )

    add_band(s, 1.4, "라우트 / 모드")
    add_kv_table(
        s,
        0.5,
        1.85,
        12.3,
        ["엔드포인트", "default 모드", "옵션 모드"],
        [
            [
                "POST /simulate-abm",
                "use_policy=True → LLM 0회",
                "enable_llm_decisions: Tier S 만 LLM",
            ],
            [
                "n_agents 파라미터",
                "5,000 (production)",
                "enable_llm_thought: Tier S thought feed 별도 LLM",
            ],
        ],
    )

    add_band(s, 3.4, "현재 활성 메모리 / 보조 시스템")
    add_kv_table(
        s,
        0.5,
        3.85,
        12.3,
        ["시스템", "역할", "활성 여부"],
        [
            ["MemoryStore", "raw 행동 로그 (deque 200) + 일일 요약 (LLM 0)", "✅ 사용"],
            ["policy_cache.json", "역할×날씨×시간 정책 (Ollama 로 11회 cold gen)", "✅ 사용"],
            ["Archetype 30+", "role 안 다양성 (ARCHETYPES dict multiplier)", "✅ 사용"],
            ["Memory Seeder", "14일 가상 visit prefill (cold start 완화)", "❌ /simulate-abm 비활성"],
            ["PgVectorMemory", "Tier S semantic search (BGE-M3)", "❌ 옵션 (default off)"],
        ],
    )

    add_band(s, 6.0, "한 줄 요약")
    add_text(
        s,
        0.5,
        6.4,
        12.3,
        0.5,
        "default = 전 Tier policy_decide (LLM 0회). Tier S 50명만 옵션으로 LLM 호출 가능.",
        size=14,
        color=GREEN,
    )


def build_db(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "DB / 데이터 (현재)", size=28, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "PostgreSQL 87 public 테이블 · ORM 77 클래스 · alembic head a8f3d2e7c1b9",
        size=13,
        color=MUTED,
    )

    add_band(s, 1.4, "핵심 테이블 (실 row count)")
    add_kv_table(
        s,
        0.5,
        1.85,
        6.0,
        ["테이블", "row 수"],
        [
            ["ftc_brand_franchise", "34,708"],
            ["biz_brand_mapping", "5,867"],
            ["users", "23 (superadmin 1)"],
            ["simulation_ai (이력)", "7"],
            ["simulation_foresee (이력)", "8"],
        ],
    )

    add_text(s, 6.8, 1.55, 6, 0.4, "외부 API 연결 (실 적재)", size=14, bold=True, color=ACCENT_2)
    add_bullets(
        s,
        6.8,
        2.0,
        6.4,
        3.0,
        [
            "서울열린데이터 — flpop / golmok / district_sales / 지하철",
            "SGIS — 인구 / 가구 / 사업체",
            "MOLIT — 부동산 실거래",
            "공정거래위원회 (FTC) — 프랜차이즈 정보공개서",
            "ECOS — 한국은행 경기 지표",
            "기상청 — weather_daily",
            "Kakao — kakao_store / kakao_store_menu",
            "공공자전거 — ttareungi.dong_code 마포 100%",
        ],
        size=11,
    )

    add_band(s, 5.0, "권한 시스템")
    add_kv_table(
        s,
        0.5,
        5.4,
        12.3,
        ["역할", "테이블", "데이터 가시 범위"],
        [
            ["master", "users", "본인 + 소속 매니저"],
            ["manager", "manager_users", "본인만"],
            ["superadmin", "users.is_superadmin=true", "전체 가맹본부 (시뮬 이력)"],
        ],
    )


def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title(prs)
    build_abm(prs)
    build_db(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "spotter-status-briefing.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"saved: {out}")
    return out


if __name__ == "__main__":
    main()
