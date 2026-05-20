"""SPOTTER 중간발표 PPT — ABM(17~20) + DB(7) 시나리오용.

이팀장 스토리 톤. 슬라이드 번호는 25장 시나리오 기준.
코드 검토 기반 정확한 수치:
- agents.py: Layer 1~5 인간다움
- vacancy_pse.py: 1일 × 5 seed, 분기/연 환산
- vacancy_inject.py: measure_cannibalization (with/without 같은 seed)
- config.py: TierDistribution 5%/20%/75%, PopulationMix v12
- models.py: 26+ 테이블 (운영 핵심)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "presentation" / "abm-midterm-2026-05-07.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GRAY = RGBColor(0x37, 0x41, 0x51)
SUB = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PINK = RGBColor(0xEC, 0x48, 0x99)
SOFT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
SOFT_RED = RGBColor(0xFE, 0xE2, 0xE2)


def add_title(slide, num, text, sub=None):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.95), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = Inches(0)
    btf.margin_right = Inches(0)
    btf.margin_top = Inches(0)
    btf.margin_bottom = Inches(0)
    bp = btf.paragraphs[0]
    bp.text = str(num)
    bp.alignment = 2
    bp.font.size = Pt(22)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(1.65), Inches(0.45), Inches(11), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if sub:
        sb = slide.shapes.add_textbox(Inches(1.65), Inches(1.05), Inches(11), Inches(0.4))
        sp = sb.text_frame.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(14)
        sp.font.italic = True
        sp.font.color.rgb = SUB

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def hero(slide, left, top, w, big, label, color=ACCENT, big_size=44):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color
    box.line.width = Pt(2.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = big
    p.alignment = 2
    p.font.size = Pt(big_size)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = 2
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY


def card(slide, left, top, w, h, title, body, color=ACCENT, title_size=20, body_size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(10)
    for item in body:
        pi = tf.add_paragraph()
        pi.text = f"· {item}"
        pi.font.size = Pt(body_size)
        pi.font.color.rgb = GRAY
        pi.space_after = Pt(5)


# ============================================================
# [7] 6 채널 → 26 테이블, 분산된 데이터를 한 곳으로
# ============================================================
def slide_07_db(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 7, "6 채널 → 26 테이블, 분산된 데이터를 한 곳으로", "PostgreSQL 메인 RDB · PGVector 법률 RAG")

    # 6 채널 (위)
    sources = [
        ("KT", "생활인구", ACCENT),
        ("카카오", "매장 16K", PURPLE),
        ("통계청", "SGIS 인구·가구", GREEN),
        ("네이버", "트렌드", ORANGE),
        ("공정위", "FTC 프랜차이즈", PINK),
        ("서울", "상권 매출·임대", RED),
    ]
    for i, (name, role, color) in enumerate(sources):
        x = 0.5 + i * 2.13
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(2.0), Inches(0.95))
        b.fill.solid()
        b.fill.fore_color.rgb = color
        b.line.fill.background()
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = 2
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = role
        p2.alignment = 2
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT

    # 화살표
    for x in (1.5, 3.6, 5.7, 7.8, 9.9, 12.0):
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x - 0.55), Inches(2.95), Inches(0.4), Inches(0.4))
        arr.fill.solid()
        arr.fill.fore_color.rgb = NAVY
        arr.line.fill.background()

    # 가운데 RDB
    rdb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.55), Inches(8.3), Inches(1.3))
    rdb.fill.solid()
    rdb.fill.fore_color.rgb = NAVY
    rdb.line.fill.background()
    tf = rdb.text_frame
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "PostgreSQL — 26 테이블"
    p.alignment = 2
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "SQLAlchemy 2.0 async · pool 25 · alembic 30+ migration"
    p2.alignment = 2
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(4)

    # 아래 분류 + PGVector
    cats = [
        ("상권/매출", "district_sales · golmok_*", ACCENT),
        ("인구", "living_pop · sgis_* · resident_*", GREEN),
        ("매장/브랜드", "kakao_store · ftc_* · biz_*", PURPLE),
        ("부동산", "rent_cost · jeonse_* · apt_*", ORANGE),
    ]
    for i, (cat, ex, color) in enumerate(cats):
        x = 0.5 + i * 3.2
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.1), Inches(3.0), Inches(1.1))
        b.fill.solid()
        b.fill.fore_color.rgb = LIGHT
        b.line.color.rgb = color
        b.line.width = Pt(1.5)
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = ex
        p2.font.size = Pt(10)
        p2.font.name = "Consolas"
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(3)

    # PGVector 별도 박스 (오른쪽)
    pg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.85))
    pg.fill.solid()
    pg.fill.fore_color.rgb = WHITE
    pg.line.color.rgb = PURPLE
    pg.line.width = Pt(2)
    tf = pg.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "+  PGVector (법률 RAG)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p2 = tf.add_paragraph()
    p2.text = "BAAI/bge-m3 1024차원 임베딩 · HNSW index · 10K+ 법률 chunk"
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(2)


# ============================================================
# [17] 카페 차리기 전, 가상으로 90일 운영해본다
# ============================================================
def slide_17_hook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 17, "카페 차리기 전, 1일 시뮬 × 5 seed로 미리 답을 본다")

    # 큰 인용
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5))
    quote.fill.solid()
    quote.fill.fore_color.rgb = LIGHT
    quote.line.color.rgb = ACCENT
    quote.line.width = Pt(3)
    tf = quote.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "“망원동 빈 가게에 카페 차리려는데, 잘 될까요?”"
    p.alignment = 2
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 좌우 비교
    card(
        slide,
        0.6,
        3.7,
        6.0,
        3.3,
        "기존 방식",
        ["부동산 수수료 + 컨설팅 수백만원", "6개월 운영 후 망하면 손해 1억", "실패 비용 = 실제 손실"],
        color=RED,
        title_size=22,
        body_size=18,
    )

    card(
        slide,
        6.85,
        3.7,
        6.0,
        3.3,
        "ABM",
        ["가상 시민 5,000명에게 카페 차려봄", "1일 시뮬 × 5 seed (PSE) → 분기/연 환산", "실패 비용 = 0원"],
        color=GREEN,
        title_size=22,
        body_size=18,
    )


# ============================================================
# [18] 5,000명, NVIDIA 7,187 페르소나로 채운다
# ============================================================
def slide_18_tier(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        18,
        "5,000명, NVIDIA 7,187 페르소나로 채운다",
        "기본 운영 LLM 0회 — Tier S Plan은 옵션 (Stanford UIST'23, default off)",
    )

    # Hero 3개
    hero(slide, 0.6, 1.85, 4.0, "7,187", "NVIDIA Nemotron 합성 페르소나", PURPLE, big_size=44)
    hero(slide, 4.7, 1.85, 4.0, "26 컬럼", "직업·학력·취미·가치관 등", ACCENT, big_size=36)
    hero(slide, 8.8, 1.85, 4.05, "0회", "vacancy 평가 LLM 호출", GREEN, big_size=44)

    # 흐름 — 7187 → sample → 5000
    flow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.1), Inches(12.15), Inches(1.4))
    flow.fill.solid()
    flow.fill.fore_color.rgb = NAVY
    flow.line.fill.background()
    tf = flow.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "Nemotron 7,187 (마포 합성)   →   sex + age_bucket 매칭 sample   →   5,000 가상 시민"
    p.alignment = 2
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "occupation · education · hobbies · professional / sports / arts / culinary persona"
    p2.alignment = 2
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(8)

    # 두뇌 — 단순화
    card(
        slide,
        0.6,
        5.7,
        6.0,
        1.5,
        "행동 결정 (default)",
        ["policy_executor.py — 순수 Python 가중합", "거리·선호·dong_affinity·시간×동×연령"],
        color=ORANGE,
        title_size=16,
        body_size=14,
    )
    card(
        slide,
        6.85,
        5.7,
        5.9,
        1.5,
        "옵션 LLM (use_llm_decisions=True)",
        ["Tier S 250명 하루 1회 plan batch (Stanford UIST'23)", "vacancy 평가는 mock 강제 → Optuna 튜닝 가능"],
        color=GREEN,
        title_size=16,
        body_size=14,
    )
    return  # skip old tier rendering

    # (legacy below, unreachable)
    for i, (tag, n, model, model2, role, color) in enumerate([]):
        x = 0.5 + i * 4.25
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.85), Inches(4.0), Inches(0.85))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        tf = head.text_frame
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = f"Tier {tag}  ·  {n}"
        p.alignment = 2
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.7), Inches(4.0), Inches(4.3))
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT
        body.line.color.rgb = color
        body.line.width = Pt(1.5)
        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = model
        p.alignment = 2
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = model2
        p2.alignment = 2
        p2.font.size = Pt(12)
        p2.font.name = "Consolas"
        p2.font.color.rgb = SUB
        p2.space_before = Pt(4)

        # 구분선
        p3 = tf.add_paragraph()
        p3.text = "─────"
        p3.alignment = 2
        p3.font.size = Pt(12)
        p3.font.color.rgb = LIGHT
        p3.space_before = Pt(20)

        for line in role.split("\n"):
            pi = tf.add_paragraph()
            pi.text = line
            pi.alignment = 2
            pi.font.size = Pt(15)
            pi.font.color.rgb = GRAY
            pi.space_after = Pt(2)


# ============================================================
# [19] 사람처럼 행동하는 가상 시민 — 5 Layer 인간다움
# ============================================================
def slide_19_layers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 19, "사람처럼 행동하는 가상 시민 — 5 Layer 인간다움")

    layers = [
        ("L1", "정책", "성향 (mobility, 카테고리 선호)\nrole × weather × time_block", ACCENT),
        ("L2", "기억", "visit_history · learned_prefs\nblacklist · habit_store", GREEN),
        ("L3", "내부 상태", "hunger · fatigue · mood\n매 tick 변화 → 의사결정 가중", ORANGE),
        ("L4", "시간 · 날씨 적응", "5 시간대 곱셈 테이블\n동 character cat_boost", PURPLE),
        ("L5", "소셜", "Dunbar k=5 친구\n만족도 > 0.7 → 추천 전파", PINK),
    ]
    for i, (tag, title, body, color) in enumerate(layers):
        x = 0.4 + i * 2.55

        # 번호 뱃지
        b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.85), Inches(1.1), Inches(1.1))
        b.fill.solid()
        b.fill.fore_color.rgb = color
        b.line.fill.background()
        tf = b.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_right = Inches(0)
        p = tf.paragraphs[0]
        p.text = tag
        p.alignment = 2
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 카드
        cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.15), Inches(2.5), Inches(3.6))
        cb.fill.solid()
        cb.fill.fore_color.rgb = WHITE
        cb.line.color.rgb = color
        cb.line.width = Pt(2)
        tf = cb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = 2
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(15)
        for line in body.split("\n"):
            pi = tf.add_paragraph()
            pi.text = line
            pi.alignment = 2
            pi.font.size = Pt(12)
            pi.font.color.rgb = GRAY
            pi.space_after = Pt(4)

    # 하단 한 줄
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4))
    p = note.text_frame.paragraphs[0]
    p.text = "agents.py · 정책만 있으면 똑같이 행동 → 5 Layer 더해서 사람마다 달라진다"
    p.alignment = 1
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = SUB


# ============================================================
# [20] 신규 진입은 시장을 키우는가, 빼앗는가
# ============================================================
def slide_20_cannibal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 20, "신규 진입은 시장을 키우는가, 빼앗는가")

    # 측정 방법 박스
    method = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.4))
    method.fill.solid()
    method.fill.fore_color.rgb = NAVY
    method.line.fill.background()
    tf = method.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "측정 방법 — 같은 seed, 두 번 시뮬"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "  baseline (vacancy 없이)   vs   with_vacancy   →   동 카페 합계 차이 = 카니발 / 시너지"
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(8)

    # 결과 hero 2개
    hero(slide, 1.0, 3.55, 5.5, "+1.4 ± 3.5%", "동 카페 시장 성장률 (95% CI)", ACCENT, big_size=38)
    hero(slide, 6.85, 3.55, 5.5, "CI 0 포함", "통계적으로 시장 확장 없음", RED, big_size=32)

    # 결론
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_GREEN
    box.line.color.rgb = GREEN
    box.line.width = Pt(2.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "발견 — zero-sum 구조"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    p2.text = "마포 카페 시장은 신규 진입에 zero-sum — 새 카페 손님은 동 다른 카페 손님 이동"
    p2.font.size = Pt(15)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(8)


# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_07_db(prs)
    slide_17_hook(prs)
    slide_18_tier(prs)
    slide_19_layers(prs)
    slide_20_cannibal(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
