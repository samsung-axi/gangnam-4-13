"""SPOTTER A1 (찬영) 작업 PPT 생성 — ABM + DB.

산출: docs/presentation/spotter-abm-db.pptx

실 코드 / DB 검증 기반 (2026-05-06):
- ABM 5,000 agents (n_agents 파라미터, /simulate-abm)
- 4가지 mode (default policy / enable_llm_decisions / 풀 JSON / DSL)
- DB: 87 public 테이블 / 77 ORM / 37 alembic migration

실행: python -m scripts.build_abm_db_pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# 색
# ---------------------------------------------------------------------------
BG = RGBColor(0x10, 0x14, 0x20)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)  # cyan
ACCENT_2 = RGBColor(0xFF, 0xB7, 0x4D)  # amber
TEXT = RGBColor(0xE3, 0xE9, 0xF4)
MUTED = RGBColor(0x9C, 0xA9, 0xBE)
GREEN = RGBColor(0x6E, 0xE7, 0xB7)


# ---------------------------------------------------------------------------
# 헬퍼
# ---------------------------------------------------------------------------
def add_dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return box


def add_bullets(slide, x, y, w, h, lines, *, size=14, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, c = item
        else:
            txt, c = item, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"·  {txt}"
        run.font.size = Pt(size)
        run.font.color.rgb = c
        run.font.name = "Pretendard"


def add_kv_table(slide, x, y, w, headers, rows, *, header_size=12, body_size=11):
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.4 * (len(rows) + 1) + 0.1),
    ).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = ACCENT
                r.font.name = "Pretendard"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x22, 0x36)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = TEXT
                    run.font.name = "Pretendard"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x28)


def add_section_band(slide, y, label):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.04))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    add_text(slide, 0.5, y + 0.05, 12, 0.4, label, size=14, bold=True, color=ACCENT)


# ---------------------------------------------------------------------------
# 슬라이드
# ---------------------------------------------------------------------------
def build_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.7, 1.6, 12, 1.0, "SPOTTER — ABM × DB 인프라", size=44, bold=True)
    add_text(s, 0.7, 2.7, 12, 0.6, "찬영 (A1) 기여 분석 · 2026-04 ~ 2026-05", size=20, color=MUTED)
    add_section_band(s, 4.0, "발표 범위")
    add_bullets(
        s,
        0.7,
        4.5,
        12,
        2.5,
        [
            "Slide 2 — ABM: 5,000 에이전트 × 4 mode × Policy-as-code",
            "Slide 3 — DB 인프라: 87 public 테이블 / 77 ORM / 37 alembic",
            "Slide 4 — 정량 가치 + 산출물",
        ],
        size=16,
    )


def build_abm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "ABM — Agent-Based Modeling", size=30, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        '마포구 5,000 가상 에이전트가 매 시간 "어디 갈지" 결정 → 매장별 매출 추정',
        size=14,
        color=MUTED,
    )

    # 4 mode 표 — 핵심 정정 (Tier A 는 default LLM 호출 X)
    add_section_band(s, 1.5, "4가지 의사결정 Mode (agents.py:decide 라우터)")
    add_kv_table(
        s,
        0.5,
        2.0,
        12.3,
        ["Mode", "Tier S (50)", "Tier A (200)", "Tier B (≈4,750)", "LLM 호출"],
        [
            ["default (use_policy=True)", "policy", "policy", "policy", "0"],
            ["enable_llm_decisions", "smart_decide (LLM)", "policy", "policy", "Tier S 만"],
            ["풀 JSON 모드", "smart_decide", "fast_decide (LLM)", "rule (0)", "S+A"],
            ["DSL 모드 (use_dsl=True)", "dsl_decide", "dsl_decide", "dsl_decide", "전 Tier"],
        ],
    )

    # 핵심 트릭
    add_section_band(s, 4.3, "핵심 설계 (LLM 비용 최소화)")
    add_bullets(
        s,
        0.5,
        4.7,
        6.4,
        2.5,
        [
            ("policy_generator — Ollama Qwen2.5:3b 로 정책 11회 생성", TEXT),
            ("  → policy_cache.json 캐시 → 매 시뮬 LLM 0회", GREEN),
            ("Archetype 30+ multiplier (resident 7종 등 role × 유형)", TEXT),
            ("  → 같은 role 이어도 행동 패턴 다양성 확보", MUTED),
            ("Memory Seeder — 14일 가상 visit_history prefill", TEXT),
            ("  → Layer 2 학습 cold start 완화 (LLM 0)", MUTED),
            ("Tier S: smart_decide (배치 LLM, Hierarchical UIST'23)", TEXT),
            ("  → /simulate-abm 전 Tier OpenAI gpt-4.1-mini 통일", MUTED),
        ],
        size=11,
    )

    # 모델
    add_text(s, 7.0, 4.7, 6, 0.4, "모델 / 인구 (실측)", size=14, bold=True, color=ACCENT_2)
    add_kv_table(
        s,
        7.0,
        5.1,
        6.0,
        ["항목", "값"],
        [
            ["LLM provider", "OpenAI gpt-4.1-mini"],
            ["Population mix", "60/20/10/4/5/1 % (resident → owner)"],
            ["base 합계 (n_personas 비례 scale)", "500"],
            ["Tier S 모드", "enable_llm_decisions=True 시"],
            ["Thought feed (시각화)", "enable_llm_thought=True (별도 LLM)"],
        ],
    )


def build_db(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "DB 인프라 / ORM 정합", size=30, bold=True, color=ACCENT)
    add_text(
        s,
        0.5,
        0.95,
        12,
        0.4,
        "PostgreSQL 87 public 테이블 · SQLAlchemy ORM 77 클래스 · Alembic 37 마이그",
        size=14,
        color=MUTED,
    )

    add_section_band(s, 1.5, "내가 한 작업 — 6 영역")

    add_bullets(
        s,
        0.5,
        2.0,
        6.4,
        4.6,
        [
            ("dong_code FK 정합화 — 4 sprint (그룹 A / B-1 / B-2 / B-3)", TEXT),
            ("  · seoul_dong_master 425동 + jeonse_dong_master 399동", MUTED),
            ("  · dong_code 8자리 검증 (helper + Pydantic validator)", MUTED),
            ("  · dong_centroid + master ORM 클래스 3종", MUTED),
            ("신규 테이블 / ORM 추가", TEXT),
            ("  · kakao_store_menu 신설 + panel3 전수 재크롤링", MUTED),
            ("  · seoul_realtime_hotspots, elderly_ratio_region", MUTED),
            ("  · emerging-trend B1 — 5 tables (master+operational)", MUTED),
            ("  · industry_master FK 배선 + InviteCode nullable 명시", MUTED),
            ("Alembic 마이그레이션 정합", TEXT),
            ("  · phantom revision 복구 + simulation_history 생성/cleanup", GREEN),
            ("  · users / manager_users 라이프사이클 컬럼 catchup", MUTED),
            ("  · users.is_superadmin 컬럼 추가 (권한 시스템)", MUTED),
        ],
        size=11,
    )

    add_bullets(
        s,
        7.0,
        2.0,
        6.0,
        4.6,
        [
            ("연결 풀 / 인프라 안정화", TEXT),
            ("  · vector_db PGVector pool_recycle=1800 (idle 회수)", MUTED),
            ("  · services 레이어 Engine 싱글톤화", GREEN),
            ("  → RDS 커넥션 포화 (max 191) 해소", GREEN),
            ("데이터 보강 (NULL / orphan)", TEXT),
            ("  · 87 테이블 매핑 전수감사 + zombie 정리", MUTED),
            ("  · NULL/orphan 감사 + master 메타 backfill", MUTED),
            ("  · 외부 API NULL fill (subway / ttareungi / hotspots)", MUTED),
            ("  · weather_daily.snow 100% 채움", MUTED),
            ("  · ETL 재적재: ttareungi.dong_code 마포 + ecos.cycle", GREEN),
            ("권한 시스템 신설", TEXT),
            ("  · users.is_superadmin BOOLEAN 컬럼 + alembic", MUTED),
            ("  · seed_superadmin.py CLI (자동 부여 금지)", MUTED),
            ("  · /admin/brands 라우터 (FTC 16K + biz_brand UNION)", MUTED),
        ],
        size=11,
    )


def build_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    add_text(s, 0.5, 0.3, 12, 0.7, "정량 가치 + 산출물", size=30, bold=True, color=ACCENT)

    add_text(s, 0.5, 1.2, 6, 0.4, "ABM 정량", size=16, bold=True, color=ACCENT_2)
    add_kv_table(
        s,
        0.5,
        1.7,
        6.0,
        ["지표", "Before", "After"],
        [
            ["LLM 호출 (5K agents × 24h, default)", "120,000", "0 (policy_cache)"],
            ["LLM 호출 (Tier S 모드)", "120,000", "1,200 (S 50명만)"],
            ["시뮬 latency (score_stores_batch)", "기준", "-49%"],
            ["LLM 토큰 (caveman prompt 압축)", "기준", "-60%"],
            ["시각화 규모", "1,000", "5,000 agents"],
        ],
    )

    add_text(s, 7.0, 1.2, 6, 0.4, "DB 정량", size=16, bold=True, color=ACCENT_2)
    add_kv_table(
        s,
        7.0,
        1.7,
        6.0,
        ["지표", "Before", "After"],
        [
            ["dong_code FK 그룹", "0", "4 그룹 (425+399 master)"],
            ["ORM ↔ DB drift 감사", "미수행", "87 테이블 정합"],
            ["RDS 커넥션 포화 (max 191)", "빈번", "싱글톤 + pool_recycle"],
            ["alembic head", "phantom", "정합 (37 마이그)"],
            ["weather/ttareungi/ecos NULL", "부분", "100%"],
        ],
    )

    add_section_band(s, 4.7, "산출물")
    add_bullets(
        s,
        0.5,
        5.1,
        12,
        2.0,
        [
            ("커밋 50+ (non-merge), PR 19+ (IM3-241/242/243/261 sprint 등)", TEXT),
            ("ABM 모듈: brain.py 1,552 LOC · runner.py 1,640 LOC · policy_executor 1,243 LOC", MUTED),
            ("DB 모듈: models.py 2,010 LOC · 37 alembic · vector_db 138 LOC", MUTED),
            ("문서: docs/issues/2026-05-05-codebase-ultrareview.md (392줄)", MUTED),
            ("도구: seed_superadmin.py / RAG trace JSONL / audit_v4 4 CV", MUTED),
        ],
        size=13,
    )


# ---------------------------------------------------------------------------
def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title(prs)
    build_abm(prs)
    build_db(prs)
    build_value(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "spotter-abm-db-v2.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"saved: {out}")
    return out


if __name__ == "__main__":
    main()
